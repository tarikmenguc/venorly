"""
Pitch Deck Uretici — Venorly
report_json'dan 10 slide yatirimci sunumu olusturur.
Dis bagimlilik: python-pptx (requirements.txt'te mevcut)

Slide yapisi (standart investor deck):
  1. Kapak
  2. Problem
  3. Cozum
  4. Pazar Buyuklugu
  5. Rekabet
  6. Is Modeli
  7. Teknik Stack
  8. GTM & ICP
  9. Finansal Ozet
 10. Cagri
"""

import io
import re
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── Renk sabitleri ─────────────────────────────────────────────────────────────
BRAND   = RGBColor(0x8B, 0x5C, 0xF6)   # mor
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1E, 0x1E, 0x1E)
MEDIUM  = RGBColor(0x50, 0x50, 0x50)
LIGHT   = RGBColor(0xF0, 0xF0, 0xF5)
GREEN   = RGBColor(0x22, 0xC5, 0x5E)
YELLOW  = RGBColor(0xEA, 0xB3, 0x08)
RED     = RGBColor(0xEF, 0x44, 0x44)

SLIDE_W = Inches(13.33)   # widescreen 16:9
SLIDE_H = Inches(7.5)


# ── Yardimci fonksiyonlar ──────────────────────────────────────────────────────

def _clean(text: Optional[str], max_len: int = 300) -> str:
    if not text:
        return ""
    text = str(text)
    text = re.sub(r'[\U00010000-\U0010ffff]', '', text)
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
    return text.strip()[:max_len]


def _bg(slide, color: RGBColor = DARK) -> None:
    """Slide arka planini doldur."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _text_box(
    slide, text: str, left, top, width, height,
    font_size=18, bold=False, color=WHITE,
    align=PP_ALIGN.LEFT, word_wrap=True,
) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _accent_bar(slide, color: RGBColor = BRAND) -> None:
    """Sol kenar aksent cubugu."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), Inches(0.12), SLIDE_H,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def _slide_header(slide, title: str, subtitle: str = "") -> None:
    _accent_bar(slide)
    _text_box(slide, title,
              left=Inches(0.3), top=Inches(0.25),
              width=Inches(12.5), height=Inches(0.8),
              font_size=28, bold=True, color=WHITE)
    if subtitle:
        _text_box(slide, subtitle,
                  left=Inches(0.3), top=Inches(0.95),
                  width=Inches(12.5), height=Inches(0.45),
                  font_size=14, color=RGBColor(0xCC, 0xCC, 0xCC))


def _bullet_list(slide, items: list, left, top, width, height, font_size=16) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "•  " + _clean(item, 120)
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)


def _stat_card(slide, label: str, value: str, left, top, width=Inches(3.5)) -> None:
    """Kucuk istatistik karti."""
    height = Inches(1.2)
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x3A)
    box.line.color.rgb = BRAND

    _text_box(slide, value, left + Inches(0.15), top + Inches(0.05),
              width - Inches(0.3), Inches(0.6),
              font_size=22, bold=True, color=WHITE)
    _text_box(slide, label, left + Inches(0.15), top + Inches(0.65),
              width - Inches(0.3), Inches(0.45),
              font_size=11, color=RGBColor(0xAA, 0xAA, 0xAA))


# ── Slide fabrikasi ───────────────────────────────────────────────────────────

def _slide_cover(prs: Presentation, rj: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _bg(slide, RGBColor(0x0F, 0x0F, 0x1A))

    # Sol cizgi kalin
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.5), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND
    bar.line.fill.background()

    title = _clean(rj.get("idea_title", "Startup Fikri"), 80)
    _text_box(slide, title,
              left=Inches(0.8), top=Inches(2.2),
              width=Inches(11), height=Inches(1.5),
              font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    es = rj.get("executive_summary") or {}
    decision = es.get("decision", "")
    score = es.get("weighted_score")
    badge_text = f"{decision}  |  Skor: {score:.0f}/100" if score else decision
    dec_color = GREEN if decision == "Go" else (YELLOW if decision == "Hold" else RED)
    _text_box(slide, badge_text,
              left=Inches(0.8), top=Inches(3.9),
              width=Inches(6), height=Inches(0.6),
              font_size=18, bold=True, color=dec_color)

    _text_box(slide, "Venorly — AI Startup Fizibilite Platformu",
              left=Inches(0.8), top=Inches(6.8),
              width=Inches(11), height=Inches(0.4),
              font_size=11, color=RGBColor(0x66, 0x66, 0x88))


def _slide_problem(prs, rj: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _slide_header(slide, "Problem", "Cozmek istedigimiz acik nedir?")

    comp = rj.get("competition") or {}
    es   = rj.get("executive_summary") or {}

    problems = []
    gap = _clean(comp.get("gap_summary", ""), 200)
    if gap:
        problems.append(gap)
    for lof in (es.get("leap_of_faith") or [])[:2]:
        problems.append(_clean(lof, 100))
    if not problems:
        problems = ["Mevcut cozumler yetersiz kaliyor", "Manuel surecler verimlilik kaybi yaratıyor"]

    _bullet_list(slide, problems[:4],
                 left=Inches(0.4), top=Inches(1.6),
                 width=Inches(12), height=Inches(4.5))


def _slide_solution(prs, rj: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _slide_header(slide, "Cozum", "Ne inşa ediyoruz?")

    v = rj.get("validation") or {}
    t = rj.get("technical") or {}

    value_prop = _clean(v.get("value_prop", ""), 200)
    stack = _clean(t.get("stack", ""), 150)

    if value_prop:
        _text_box(slide, value_prop,
                  left=Inches(0.4), top=Inches(1.6),
                  width=Inches(12), height=Inches(1.5),
                  font_size=20, color=WHITE)

    items = []
    for step in (v.get("cold_email_sequence") or [])[:3]:
        items.append(_clean(step, 100))
    if not items and stack:
        items = [f"Teknoloji: {stack}"]

    if items:
        _bullet_list(slide, items,
                     left=Inches(0.4), top=Inches(3.3),
                     width=Inches(12), height=Inches(3))


def _slide_market(prs, rj: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _slide_header(slide, "Pazar Buyuklugu", "Adreslenebilir firsatin boyutu")

    m = rj.get("market") or {}
    y = Inches(1.7)
    spacing = Inches(1.4)

    for label, key in [("TAM", "tam"), ("SAM", "sam"), ("SOM", "som")]:
        val = _clean(m.get(key, "—"), 40)
        formula = _clean(m.get(f"{key}_formula", ""), 80) if key == "tam" else ""
        _stat_card(slide, label + (f"\n{formula}" if formula else ""), val,
                   left=Inches(0.4), top=y)
        y += spacing

    cagr = _clean(m.get("cagr", ""), 30)
    if cagr:
        _stat_card(slide, "CAGR (Yillik Buyume)", cagr,
                   left=Inches(4.2), top=Inches(1.7))

    tam_src = _clean(m.get("tam_source", ""), 60)
    if tam_src:
        _text_box(slide, f"Kaynak: {tam_src}",
                  left=Inches(0.4), top=Inches(6.7),
                  width=Inches(12), height=Inches(0.4),
                  font_size=10, color=RGBColor(0x88, 0x88, 0x88))


def _slide_competition(prs, rj: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _slide_header(slide, "Rekabet", "Kim var, ne eksik?")

    c = rj.get("competition") or {}
    competitors = (c.get("competitors") or [])[:4]

    y = Inches(1.6)
    for comp in competitors:
        name    = _clean(comp.get("name", ""), 30)
        weakness = _clean(comp.get("weakness", "Veri yok"), 80)
        _text_box(slide, f"{name}: {weakness}",
                  left=Inches(0.4), top=y,
                  width=Inches(12), height=Inches(0.55),
                  font_size=14, color=RGBColor(0xDD, 0xDD, 0xDD))
        y += Inches(0.65)

    gap = _clean(c.get("gap_summary", ""), 200)
    if gap:
        _text_box(slide, "Bizim farkımız: " + gap,
                  left=Inches(0.4), top=max(y, Inches(5.0)),
                  width=Inches(12), height=Inches(1.2),
                  font_size=15, bold=True, color=BRAND)


def _slide_business_model(prs, rj: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _slide_header(slide, "Is Modeli", "Nasil para kazaniyoruz?")

    t = rj.get("technical") or {}
    pricing = _clean(t.get("pricing_model", ""), 150)
    ltv     = _clean(t.get("ltv", ""), 40)
    cac     = _clean(t.get("cac", ""), 40)
    cpu     = _clean(t.get("cpu_cost", ""), 40)

    if pricing:
        _text_box(slide, pricing,
                  left=Inches(0.4), top=Inches(1.6),
                  width=Inches(12), height=Inches(1.0),
                  font_size=18, color=WHITE)

    x = Inches(0.4)
    for label, val in [("LTV", ltv), ("CAC", cac), ("CPU Maliyeti", cpu)]:
        if val:
            _stat_card(slide, label, val, left=x, top=Inches(3.0))
            x += Inches(4.0)


def _slide_tech_stack(prs, rj: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _slide_header(slide, "Teknik Altyapi", "Nasil insaediyoruz?")

    t = rj.get("technical") or {}
    stack = _clean(t.get("stack", ""), 200)

    if stack:
        _text_box(slide, stack,
                  left=Inches(0.4), top=Inches(1.7),
                  width=Inches(12), height=Inches(2.0),
                  font_size=18, color=WHITE)

    ensemble = (rj.get("executive_summary") or {}).get("ensemble_note", "")
    if ensemble:
        _text_box(slide, "AI Karar Konsensus: " + _clean(ensemble, 120),
                  left=Inches(0.4), top=Inches(4.5),
                  width=Inches(12), height=Inches(0.6),
                  font_size=13, color=RGBColor(0x99, 0x99, 0xCC))


def _slide_gtm(prs, rj: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _slide_header(slide, "GTM & Hedef Kitle", "Ilk musteriye nasil ulasiriz?")

    v = rj.get("validation") or {}
    icp   = _clean(v.get("icp", ""), 120)
    h1    = _clean(v.get("waitlist_h1", ""), 80)
    dm    = _clean(v.get("linkedin_dm", ""), 200)

    if icp:
        _text_box(slide, "ICP: " + icp,
                  left=Inches(0.4), top=Inches(1.6),
                  width=Inches(12), height=Inches(0.7),
                  font_size=17, bold=True, color=BRAND)
    if h1:
        _text_box(slide, h1,
                  left=Inches(0.4), top=Inches(2.5),
                  width=Inches(12), height=Inches(0.8),
                  font_size=16, color=WHITE)
    if dm:
        _text_box(slide, "Cold DM ornegi: " + dm,
                  left=Inches(0.4), top=Inches(3.5),
                  width=Inches(12), height=Inches(2.5),
                  font_size=13, color=RGBColor(0xCC, 0xCC, 0xCC))


def _slide_financials(prs, rj: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _slide_header(slide, "Finansal Ozet", "Birim ekonomisi ve projeksiyon")

    t  = rj.get("technical") or {}
    m  = rj.get("market") or {}
    es = rj.get("executive_summary") or {}

    stats = [
        ("LTV", _clean(t.get("ltv", "—"), 30)),
        ("CAC", _clean(t.get("cac", "—"), 30)),
        ("SOM (Yil 1)", _clean(m.get("som", "—"), 30)),
        ("Agirlikli Skor", f'{es.get("weighted_score", 0):.0f}/100'),
    ]
    x = Inches(0.4)
    row = 0
    for label, val in stats:
        col = row % 2
        lft = Inches(0.4 + col * 6.3)
        tp  = Inches(1.8 + (row // 2) * 1.5)
        _stat_card(slide, label, val, left=lft, top=tp, width=Inches(5.8))
        row += 1

    ci = rj.get("confidence_index")
    if ci is not None:
        ci_color = GREEN if float(ci) >= 0.75 else (YELLOW if float(ci) >= 0.5 else RED)
        _text_box(slide, f"Guven Endeksi: {float(ci):.0%}",
                  left=Inches(0.4), top=Inches(6.5),
                  width=Inches(6), height=Inches(0.5),
                  font_size=14, bold=True, color=ci_color)


def _slide_ask(prs, rj: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, RGBColor(0x0F, 0x0F, 0x1A))

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.5), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND
    bar.line.fill.background()

    _text_box(slide, "Siradaki Adim",
              left=Inches(0.8), top=Inches(1.5),
              width=Inches(11), height=Inches(1.0),
              font_size=36, bold=True, color=WHITE)

    v = rj.get("validation") or {}
    h2 = _clean(v.get("waitlist_h2", "Erken erisim icin bize ulasın."), 200)
    _text_box(slide, h2,
              left=Inches(0.8), top=Inches(2.7),
              width=Inches(11), height=Inches(1.5),
              font_size=20, color=RGBColor(0xCC, 0xCC, 0xCC))

    pivots = rj.get("pivot_suggestions")
    if pivots:
        _text_box(slide, "Alternatif Yonler:",
                  left=Inches(0.8), top=Inches(4.5),
                  width=Inches(11), height=Inches(0.5),
                  font_size=13, bold=True, color=RGBColor(0x99, 0x99, 0xCC))
        _bullet_list(slide, pivots[:3],
                     left=Inches(0.8), top=Inches(5.1),
                     width=Inches(11), height=Inches(1.8), font_size=13)

    _text_box(slide, "venorly.com",
              left=Inches(0.8), top=Inches(7.1),
              width=Inches(4), height=Inches(0.3),
              font_size=11, color=RGBColor(0x66, 0x66, 0x88))


# ── Ana fonksiyon ──────────────────────────────────────────────────────────────

def generate_pitch_deck(report_json: dict) -> bytes:
    """
    report_json (FeasibilityReport dict) → PPTX bytes.
    10 slide standart investor deck.
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_cover(prs, report_json)
    _slide_problem(prs, report_json)
    _slide_solution(prs, report_json)
    _slide_market(prs, report_json)
    _slide_competition(prs, report_json)
    _slide_business_model(prs, report_json)
    _slide_tech_stack(prs, report_json)
    _slide_gtm(prs, report_json)
    _slide_financials(prs, report_json)
    _slide_ask(prs, report_json)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
